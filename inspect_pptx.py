import sys
from pptx import Presentation

def inspect_pptx(filepath):
    prs = Presentation(filepath)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"[{shape.name}] {shape.text[:100]}")
        print()

if __name__ == "__main__":
    inspect_pptx(sys.argv[1])
